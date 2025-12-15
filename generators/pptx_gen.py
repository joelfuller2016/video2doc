"""
PPTX Generator - Creates Microsoft PowerPoint presentations from video analysis
Produces professional slide decks with screenshots and bullet points
"""

import os
import sys
from pathlib import Path
from typing import List, Dict, Optional
from datetime import datetime

# Add parent directory to path for imports when running as module
sys.path.insert(0, str(Path(__file__).parent.parent))
from logger import get_logger
from utils.path_validator import (
    is_safe_path,
    validate_file_output_path,
)

# Module logger
logger = get_logger(__name__)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def generate_pptx(
    title: str,
    summary: str,
    sections: List[Dict],
    screenshots: Dict[float, str],
    output_path: str,
    author: str = "FrameNotes"
) -> str:
    """
    Generate a Microsoft PowerPoint presentation from video analysis.

    Args:
        title: Presentation title
        summary: Brief summary for title slide
        sections: List of section dictionaries with title, content, screenshot_timestamp
        screenshots: Dict mapping timestamps to image file paths
        output_path: Path for output .pptx file
        author: Presentation author name

    Returns:
        Path to generated presentation
    """
    logger.info(f"Generating PPTX presentation: {output_path}")
    logger.debug(f"Presentation has {len(sections)} sections, {len(screenshots)} screenshots")

    # Validate output path for security
    output_result = validate_file_output_path(
        output_path,
        allowed_extensions=[".pptx"],
        create_parent_dirs=True
    )
    if not output_result.is_valid:
        raise ValueError(f"Invalid output path: {output_result.error_message}")
    output_path = output_result.sanitized_value

    # Validate screenshot paths for security
    validated_screenshots = {}
    for ts, screenshot_path in screenshots.items():
        is_safe, error = is_safe_path(screenshot_path)
        if is_safe and os.path.exists(screenshot_path):
            validated_screenshots[ts] = screenshot_path
        else:
            logger.warning(f"Skipping invalid screenshot path at {ts}: {error or 'File not found'}")
    screenshots = validated_screenshots

    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Define layouts
    title_layout = prs.slide_layouts[0]  # Title slide
    title_content_layout = prs.slide_layouts[1]  # Title and content
    blank_layout = prs.slide_layouts[6]  # Blank

    # ===== Title Slide =====
    title_slide = prs.slides.add_slide(title_layout)

    # Set title
    title_shape = title_slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Set subtitle
    subtitle_shape = title_slide.placeholders[1]
    subtitle_shape.text = summary if summary else "AI-Generated Documentation"
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    # ===== Table of Contents Slide =====
    toc_slide = prs.slides.add_slide(title_content_layout)
    toc_slide.shapes.title.text = "Contents"

    content_shape = toc_slide.placeholders[1]
    tf = content_shape.text_frame

    for i, section in enumerate(sections):
        if i == 0:
            tf.paragraphs[0].text = f"{i+1}. {section.get('title', 'Section')}"
        else:
            p = tf.add_paragraph()
            p.text = f"{i+1}. {section.get('title', 'Section')}"
            p.level = 0

    # ===== Section Slides =====
    logger.debug("Creating section slides...")
    for i, section in enumerate(sections, 1):
        section_title = section.get("title", f"Section {i}")
        section_content = section.get("content", "")
        screenshot_ts = section.get("screenshot_timestamp")
        screenshot_reason = section.get("screenshot_reason", "")

        # Find screenshot
        screenshot_path = None
        if screenshot_ts is not None:
            if screenshot_ts in screenshots:
                screenshot_path = screenshots[screenshot_ts]
            else:
                closest_ts = min(screenshots.keys(), key=lambda x: abs(x - screenshot_ts), default=None)
                if closest_ts is not None and abs(closest_ts - screenshot_ts) < 5:
                    screenshot_path = screenshots[closest_ts]

        image_added = False
        if screenshot_path and os.path.exists(screenshot_path):
            try:
                logger.debug(f"Adding slide with screenshot for section {i}: {screenshot_path}")
                # Slide with image (use blank layout for more control)
                slide = prs.slides.add_slide(blank_layout)

                # Add title at top
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3),
                    Inches(12.333), Inches(0.8)
                )
                title_tf = title_box.text_frame
                title_p = title_tf.paragraphs[0]
                title_p.text = f"{i}. {section_title}"
                title_p.font.size = Pt(32)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(0, 51, 102)

                # Add image (centered, large)
                img_left = Inches(1.0)
                img_top = Inches(1.3)
                img_width = Inches(11.333)
                img_height = Inches(5.0)

                # Calculate aspect ratio to fit
                from PIL import Image
                with Image.open(screenshot_path) as img:
                    img_ratio = img.width / img.height
                    box_ratio = img_width.inches / img_height.inches

                    if img_ratio > box_ratio:
                        # Image is wider - fit to width
                        actual_width = img_width
                        actual_height = Inches(img_width.inches / img_ratio)
                    else:
                        # Image is taller - fit to height
                        actual_height = img_height
                        actual_width = Inches(img_height.inches * img_ratio)

                    # Center the image
                    actual_left = Inches((13.333 - actual_width.inches) / 2)
                    actual_top = Inches(1.3 + (5.0 - actual_height.inches) / 2)

                slide.shapes.add_picture(
                    screenshot_path,
                    actual_left, actual_top,
                    actual_width, actual_height
                )

                # Add caption below image
                if screenshot_reason:
                    caption_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(6.5),
                        Inches(12.333), Inches(0.5)
                    )
                    caption_tf = caption_box.text_frame
                    caption_p = caption_tf.paragraphs[0]
                    caption_p.text = screenshot_reason
                    caption_p.font.size = Pt(14)
                    caption_p.font.italic = True
                    caption_p.font.color.rgb = RGBColor(100, 100, 100)
                    caption_p.alignment = PP_ALIGN.CENTER

                image_added = True

            except Exception as e:
                logger.warning(f"Failed to add image for section {i}: {e}. Creating text-only slide.")
                # Fall through to text-only slide creation

        if not image_added:
            # Slide without image (standard content layout)
            slide = prs.slides.add_slide(title_content_layout)
            slide.shapes.title.text = f"{i}. {section_title}"

            # Add content as bullet points
            if section_content:
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame

                # Split content into sentences for bullet points
                sentences = section_content.replace(". ", ".|").split("|")
                sentences = [s.strip() for s in sentences if s.strip()]

                for j, sentence in enumerate(sentences):
                    if j == 0:
                        tf.paragraphs[0].text = sentence
                        tf.paragraphs[0].font.size = Pt(20)
                    else:
                        p = tf.add_paragraph()
                        p.text = sentence
                        p.font.size = Pt(20)
                        p.level = 0

    # ===== Closing Slide =====
    end_slide = prs.slides.add_slide(blank_layout)

    # Thank you text
    thanks_box = end_slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(1.5)
    )
    thanks_tf = thanks_box.text_frame
    thanks_p = thanks_tf.paragraphs[0]
    thanks_p.text = "Thank You"
    thanks_p.font.size = Pt(54)
    thanks_p.font.bold = True
    thanks_p.font.color.rgb = RGBColor(0, 51, 102)
    thanks_p.alignment = PP_ALIGN.CENTER

    # Generation info
    info_box = end_slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(12.333), Inches(1.0)
    )
    info_tf = info_box.text_frame
    info_p = info_tf.paragraphs[0]
    info_p.text = f"Generated by FrameNotes | {datetime.now().strftime('%B %d, %Y')}"
    info_p.font.size = Pt(18)
    info_p.font.color.rgb = RGBColor(128, 128, 128)
    info_p.alignment = PP_ALIGN.CENTER

    # Ensure output directory exists
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(str(output_path))
    logger.info(f"PPTX presentation saved: {output_path}")

    return str(output_path)


if __name__ == "__main__":
    from logger import init_logging
    init_logging(verbose=True)

    # Test with mock data
    test_sections = [
        {
            "title": "Introduction",
            "content": "This presentation covers the key concepts from the video. We will explore each topic in detail.",
            "screenshot_timestamp": 5.0,
            "screenshot_reason": "Opening screen"
        },
        {
            "title": "Main Concepts",
            "content": "The main ideas are presented here. Understanding these fundamentals is essential for success.",
            "screenshot_timestamp": 30.0,
            "screenshot_reason": "Core interface"
        },
        {
            "title": "Conclusion",
            "content": "We covered all the important points. Apply these concepts to improve your workflow.",
            "screenshot_timestamp": 60.0,
            "screenshot_reason": "Final results"
        }
    ]

    output = generate_pptx(
        title="Test Presentation",
        summary="AI-Generated from Video Content",
        sections=test_sections,
        screenshots={},
        output_path="./test_output.pptx"
    )

    logger.info(f"Generated: {output}")
